# -*- coding: utf-8 -*-
"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html
Program takes a PowerPoint input file and generates a marked up version that
shows the various layouts and placeholders in the template.
"""

from __future__ import print_function
from pptx import Presentation
import argparse


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Analyze powerpoint file structure')
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file to be analyzed')
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint')
    return parser.parse_args()


def analyze_ppt(input):
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    comments from Chen: This function creates a slide for each slide layout 20190305
    """
    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            print('text of the shape is {}'.format(shape.text))
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))


def slide_analyzer(infile):
    prs = Presentation(infile)
    # load the presentation
    slide_0 = prs.slides[0]
    # load the first slide of the PPT
    for shape in slide_0.shapes:
        print(slide_0.shapes.index(shape))
        try:
            if 'text' in dir(shape):
                print(dir(shape))
                print(shape.text)
                _text_full = shape.text + '111.1, 68.8'
                if shape.is_placeholder == 1:
                    print(shape.placeholder_format.type)
                    _text_original = shape.text
                    _text_full = shape.text + '111.1, 68.8'
                    shape.text = _text_full
                    print(shape.text)
                    print(shape.shape_id)
                else:
                    print('{} is not a placeholder'.format(shape.shape_id))
        except AttributeError:
            print('This shape has no text')
    prs.save('filled_template.pptm')


if __name__ == "__main__":
#    analyze_ppt('D://工作//AUDI_L3_WP1_Virtual_Drive//PPT//checkbox-test.pptm', 'D://工作//AUDI_L3_WP1_Virtual_Drive//new-file-name.pptm')
    slide_analyzer('D://工作//AUDI_L3_WP1_Virtual_Drive//new.pptm')