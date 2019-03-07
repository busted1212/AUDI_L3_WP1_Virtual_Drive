import pptx
from pptx.util import Inches, Pt
from read_path_points import read_path_points
import six
import copy

# next job is to find how to insert check box
# and drag down menu


def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders)
                          for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def duplicate_slide(pres, index):
    """Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]
    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shape in source.shapes:
        newel = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in six.iteritems(source.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            dest.part.rels.add_relationship(value.reltype,
                                            value._target,
                                            value.rId)


prs = pptx.Presentation('D://工作//AUDI_L3_WP1_Virtual_Drive//checkbox-test.pptm')
duplicate_slide(prs, 0)  # duplicate the first slide
prs.save('output.pptm')