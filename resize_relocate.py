from PIL import Image
import pptx
from pptx.util import Inches, Pt


left = top = width = height = Inches(1.0)


def resize_and_relocate(height, width, left, top, ):
    _file_path = 'Beijing_Virtual_Drive/airport_express/ae_tocity/full_bj_ae_tocity/pic@lng=116.4641,lat=39.966324.jpg'
    _image = Image.open(_file_path)
    _image_resized = _image.resize((height, width), Image.ANTIALIAS)
    _image_resized.save('new.jpg')
    _prs = pptx.Presentation('slide_template_20190305_Chen.pptm')
    _prs.slides[0].shapes.add_picture('new.jpg', left, top)
    _prs.save('position_test_2.pptm')


if __name__ == "__main__":
    resize_and_relocate(600, 280, 0.8*left, 2.5*top)