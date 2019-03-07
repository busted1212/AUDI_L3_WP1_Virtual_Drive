# -*- coding: utf-8 -*-
from PIL import Image
import pptx
from pptx.util import Inches, Pt
from read_path_points import read_path_points

left = top = width = height = Inches(1.0)

prs = pptx.Presentation('PPT//checkbox-test.pptm')
slide_0 = prs.slides[0]
file = "D://工作//AUDI_L3_WP1_Virtual_Drive//BSV//Panorama//new.jpg"
slide_0.shapes.add_picture(file, left, top)
prs.save('new_1.pptm')


# This piece of code is able to generate ppt and insert picture automatically
def generate_ppt(points_list):
    # given a list of points, generate PPT with picture inserted
    # load a ppt template
    _prs = pptx.Presentation('slide_template_20190305_Chen.pptm')
    # select the first slide
    _slide_0 = _prs.slides[0]
    # read the layout of the fist layout
    _slide_0_layout = _slide_0.slide_layout
    # add new slide and insert contents
    for i, point in enumerate(points_list):
        # for each point
        # generate a slide to hold contents
        # specify the file to be inserted
        _file_dir = 'D://工作//AUDI_L3_WP1_Virtual_Drive//Beijing_Virtual_Drive//airport_express//ae_tocity//resized_full//'
        _file_name = 'pic@lng=' + str(point[0]) + ',lat=' + str(point[1]) + '.jpg'
        _file_path = _file_dir + _file_name
        _pic = Image.open(_file_path)
        _prs.slides[i].shapes.add_picture(_file_path, 0.75*left, 2.5*top)
        # add hyperlink
        _url_click_action = _prs.slides[i].shapes[8].click_action
        _hyperlink = _url_click_action.hyperlink
        _hyperlink.address = str(point[2])
        print('hyperlink is ' + _hyperlink.address)
        # add GPS data
        _gps = _prs.slides[i].shapes[9].text + 'lng:' + str(point[0]) + 'lat:' + str(point[1])
        _prs.slides[i].shapes[9].text = _gps
        #_prs.save('new.pptm')


# next job is to find how to insert check box
# and drag down menu
# slide_1 = prs.slides.add_slide(slide_0_layout)
# slide_0.shapes.add_picture('D://工作//AUDI_L3_WP1_Virtual_Drive//BSV//Panorama//Panoramapano@lng=116.4641lat=39.966324.jpeg', left, top)
# slide_1.shapes.add_picture('D://工作//AUDI_L3_WP1_Virtual_Drive//BSV//Panorama//Panoramapano@lng=116.4641lat=39.966324.jpeg', left, top)
# prs.save('D://工作//AUDI_L3_WP1_Virtual_Drive//Beijing_Virtual_Drive//airport_express//ae_tocity//powerpoint')


if __name__ == '__main__':
    points = read_path_points('D://工作//AUDI_L3_WP1_Virtual_Drive//Beijing_Virtual_Drive//airport_express//ae_tocity//file_list.txt')
    generate_ppt(points)
